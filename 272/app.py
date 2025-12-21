import gradio as gr
import numpy as np
import random
import torch
import tempfile
import os

from PIL import Image
from pptx import Presentation
from torchvision import transforms
from transformers import AutoModelForImageSegmentation
from diffusers import QwenImageLayeredPipeline, QwenImageEditPlusPipeline

# --- Constants ---
MAX_SEED = np.iinfo(np.int32).max
dtype = torch.bfloat16
device = "cuda" if torch.cuda.is_available() else "cpu"

# --- Model Loading ---
print("Loading Qwen-Image-Layered pipeline...")
layered_pipeline = QwenImageLayeredPipeline.from_pretrained("checkpoints/Qwen-Image-Layered")
layered_pipeline = layered_pipeline.to(device, dtype)
layered_pipeline.set_progress_bar_config(disable=None)
print("Qwen-Image-Layered pipeline loaded.")

print("Loading Qwen-Image-Edit pipeline...")
edit_pipeline = QwenImageEditPlusPipeline.from_pretrained("checkpoints/Qwen-Image-Edit-2509", torch_dtype=dtype).to(device)
print("Qwen-Image-Edit pipeline loaded.")

print("Loading RMBG model...")
rmbg_model = AutoModelForImageSegmentation.from_pretrained('checkpoints/RMBG-2.0', trust_remote_code=True).eval().to(device)
rmbg_transforms = transforms.Compose([
    transforms.Resize((1024, 1024)),
    transforms.ToTensor(),
    transforms.Normalize([0.485, 0.456, 0.406], [0.229, 0.224, 0.225])
])
print("RMBG model loaded.")


# =============================================================================
# Part 1: Image Layered Decomposition Functions
# =============================================================================

def imagelist_to_pptx(img_files):
    """Convert a list of image files to a PPTX presentation."""
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px, dpi=96):
        inch = px / dpi
        emu = inch * 914400
        return int(emu)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = 0
    for img_path in img_files:
        slide.shapes.add_picture(img_path, left, top, width=px_to_emu(img_width_px), height=px_to_emu(img_height_px))

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name


def export_gallery(images):
    """Export gallery images to a PPTX file."""
    images = [e[0] for e in images]
    pptx_path = imagelist_to_pptx(images)
    return pptx_path


def infer_decompose(
    input_image,
    seed=777,
    randomize_seed=False,
    prompt=None,
    neg_prompt=" ",
    true_guidance_scale=4.0,
    num_inference_steps=50,
    layer=4,
    cfg_norm=True,
    use_en_prompt=True
):
    """Decompose an image into multiple layers."""
    if randomize_seed:
        seed = random.randint(0, MAX_SEED)
        
    if isinstance(input_image, list):
        input_image = input_image[0]
        
    if isinstance(input_image, str):
        pil_image = Image.open(input_image).convert("RGB").convert("RGBA")
    elif isinstance(input_image, Image.Image):
        pil_image = input_image.convert("RGB").convert("RGBA")
    elif isinstance(input_image, np.ndarray):
        pil_image = Image.fromarray(input_image).convert("RGB").convert("RGBA")
    else:
        raise ValueError("Unsupported input_image type: %s" % type(input_image))
    
    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device=device).manual_seed(seed),
        "true_cfg_scale": true_guidance_scale,
        "prompt": prompt,
        "negative_prompt": neg_prompt,
        "num_inference_steps": num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": layer,
        "resolution": 640,
        "cfg_normalize": cfg_norm,
        "use_en_prompt": use_en_prompt, 
    }
    print(f"[Decompose] Running with inputs: {inputs}")
    
    with torch.inference_mode():
        output = layered_pipeline(**inputs)
        output_images = output.images[0]
    
    return output_images


# =============================================================================
# Part 2: RGBA Image Edit Functions
# =============================================================================

def blend_with_green_bg(input_img):
    """Blend an RGBA image with a green background."""
    bg = Image.new("RGB", input_img.size, (30, 215, 96)).convert("RGBA")
    input_rgba = input_img.convert("RGBA")
    blended = Image.alpha_composite(bg, input_rgba).convert("RGB")
    return blended


def infer_edit(
    image,
    prompt,
    seed=42,
    randomize_seed=False,
    true_guidance_scale=4.0,
    num_inference_steps=50,
    progress=gr.Progress(track_tqdm=True),
):
    """Edit an image using the Qwen-Image-Edit pipeline."""
    negative_prompt = " "
    
    if randomize_seed:
        seed = random.randint(0, MAX_SEED)

    generator = torch.Generator(device=device).manual_seed(seed)
    
    print(f"[Edit] Prompt: '{prompt}'")
    print(f"[Edit] Seed: {seed}, Steps: {num_inference_steps}, Guidance: {true_guidance_scale}")

    # Convert to RGB if needed
    if image.mode != "RGB":
        image = image.convert("RGB")

    # Generate the edited image
    edited_image = edit_pipeline(
        image,
        prompt=prompt,
        negative_prompt=negative_prompt,
        num_inference_steps=num_inference_steps,
        generator=generator,
        true_cfg_scale=true_guidance_scale,
        num_images_per_prompt=1
    ).images[0]

    # Remove background to create RGBA version
    input_images = rmbg_transforms(edited_image).unsqueeze(0).to(device)
    with torch.no_grad():
        preds = rmbg_model(input_images)[-1].sigmoid().cpu()
    pred = preds[0].squeeze()
    pred_pil = transforms.ToPILImage()(pred)
    mask = pred_pil.resize(edited_image.size)
    
    # Create RGBA version with transparent background
    edited_rgba = edited_image.copy()
    edited_rgba.putalpha(mask)
    
    return edited_image, edited_rgba, seed


# =============================================================================
# Examples
# =============================================================================

decompose_examples = [
    ["assets/test_images/1.png", "Illustration of three business people raising hands in celebration"],
    ["assets/test_images/2.png", "Yellow skateboard poster design with cartoon boy"],
    ["assets/test_images/3.png", "Girls night out party illustration with purple theme"],
    ["assets/test_images/4.png", "Happy Halloween greeting card with cute characters"],
    ["assets/test_images/5.png", "Happy Birthday cake celebration illustration"],
    ["assets/test_images/6.png", "Halloween party poster with purple haunted house"],
    ["assets/test_images/7.png", "Open storybook with autumn scene and birds"],
    ["assets/test_images/8.png", "Halloween cartoon scene with children and pumpkins"],
    ["assets/test_images/9.png", "Chinese anti-smoking public service announcement"],
    ["assets/test_images/10.png", "Mermaid illustration with colorful hair and ocean"],
    ["assets/test_images/11.png", "Trick or treat Halloween poster with pumpkin"],
    ["assets/test_images/12.png", "Fantasy landscape with girl and black cat"],
    ["assets/test_images/13.png", "Japanese temple scene with autumn leaves"],
]


# =============================================================================
# Gradio UI
# =============================================================================

css = """
#col-container {
    margin: 0 auto;
    max-width: 1200px;
}
"""

with gr.Blocks(css=css, title="Qwen图像分层与编辑") as demo:
    gr.Markdown("# 🎨 Qwen 图像分层与编辑")
    gr.Markdown("本应用集成了 **图像分层分解** 和 **RGBA图像编辑** 功能。")
    
    with gr.Tabs():
        # =================================================================
        # Tab 1: Image Layered Decomposition
        # =================================================================
        with gr.TabItem("🔀 图像分层分解"):
            gr.Markdown("### 将图像分解为多个图层并导出为PPTX")
            
            with gr.Row():
                with gr.Column():
                    decompose_input = gr.Image(label="输入图像", image_mode="RGBA")
                    
                with gr.Column():
                    decompose_seed = gr.Slider(
                        label="随机种子",
                        minimum=0,
                        maximum=MAX_SEED,
                        step=1,
                        value=0,
                    )
                    decompose_randomize_seed = gr.Checkbox(label="随机化种子", value=True)

                    decompose_prompt = gr.Textbox(
                        label="提示词（可选）",
                        placeholder="Enter a prompt to guide decomposition (Optional)",
                        value="",
                        lines=2,
                    )
                    decompose_neg_prompt = gr.Textbox(
                        label="负面提示词（可选）",
                        placeholder="Enter negative prompt",
                        value=" ",
                        lines=2,
                    )
                    
                    with gr.Row():
                        decompose_guidance = gr.Slider(
                            label="引导比例",
                            minimum=1.0,
                            maximum=10.0,
                            step=0.1,
                            value=4.0
                        )
                        decompose_steps = gr.Slider(
                            label="推理步数",
                            minimum=1,
                            maximum=50,
                            step=1,
                            value=50,
                        )
                        decompose_layers = gr.Slider(
                            label="图层数",
                            minimum=2,
                            maximum=10,
                            step=1,
                            value=4,
                        )

                    with gr.Row():
                        decompose_cfg_norm = gr.Checkbox(label="启用CFG归一化", value=True)
                        decompose_use_en = gr.Checkbox(label="自动描述（勾选为英文，取消为中文）", value=True)
                    
                    decompose_btn = gr.Button("🚀 开始分解!", variant="primary")

            decompose_gallery = gr.Gallery(label="分解后的图层", columns=4, rows=1, format="png")
            
            with gr.Row():
                export_btn = gr.Button("📥 导出为PPTX")
                export_file = gr.File(label="下载PPTX")
            
            export_btn.click(
                fn=export_gallery,
                inputs=decompose_gallery,
                outputs=export_file
            )
            
            gr.Examples(
                examples=decompose_examples,
                inputs=[decompose_input, decompose_prompt], 
                outputs=[decompose_gallery], 
                fn=infer_decompose, 
                examples_per_page=14,
                cache_examples=False,
                run_on_click=True
            )

            decompose_btn.click(
                fn=infer_decompose,
                inputs=[
                    decompose_input,
                    decompose_seed,
                    decompose_randomize_seed,
                    decompose_prompt,
                    decompose_neg_prompt,
                    decompose_guidance,
                    decompose_steps,
                    decompose_layers,
                    decompose_cfg_norm,
                    decompose_use_en,
                ],
                outputs=decompose_gallery,
            )

        # =================================================================
        # Tab 2: Image Editing
        # =================================================================
        with gr.TabItem("✏️ 图像编辑"):
            gr.Markdown("### AI图像编辑 + 智能去背景")
            gr.Markdown("上传任意图片，输入编辑指令进行AI编辑，同时自动生成去除背景的透明图片！")
            
            with gr.Row():
                with gr.Column():
                    edit_input = gr.Image(label="输入图像", type="pil")
                    edit_prompt = gr.Text(
                        label="编辑指令",
                        placeholder="Describe the edit instruction...",
                        container=True,
                    )
                    edit_btn = gr.Button("✨ 开始编辑!", variant="primary")

                with gr.Column():
                    edit_result = gr.Image(label="编辑结果", type="pil")
                    edit_result_rgba = gr.Image(label="去背景结果（RGBA）", type="pil")

            with gr.Accordion("高级设置", open=False):
                edit_seed = gr.Slider(
                    label="随机种子",
                    minimum=0,
                    maximum=MAX_SEED,
                    step=1,
                    value=0,
                )
                edit_randomize_seed = gr.Checkbox(label="随机化种子", value=True)

                with gr.Row():
                    edit_guidance = gr.Slider(
                        label="引导比例",
                        minimum=1.0,
                        maximum=10.0,
                        step=0.1,
                        value=4.0
                    )
                    edit_steps = gr.Slider(
                        label="推理步数",
                        minimum=1,
                        maximum=50,
                        step=1,
                        value=50,
                    )

            gr.on(
                triggers=[edit_btn.click, edit_prompt.submit],
                fn=infer_edit,
                inputs=[
                    edit_input,
                    edit_prompt,
                    edit_seed,
                    edit_randomize_seed,
                    edit_guidance,
                    edit_steps,
                ],
                outputs=[edit_result, edit_result_rgba, edit_seed],
            )


if __name__ == "__main__":
    demo.launch(
        server_name="0.0.0.0",
        server_port=7860,
    )
